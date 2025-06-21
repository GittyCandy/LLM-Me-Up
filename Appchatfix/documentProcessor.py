import sys
import PyPDF2
import textract
import docx
import pandas as pd
from pptx import Presentation
import os
from typing import Optional
import re

def extract_text_from_pdf(pdf_path: str) -> Optional[str]:
    try:
        with open(pdf_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            text = ""
            for page in reader.pages:
                text += page.extract_text() + "\n"
            return clean_text(text)
    except Exception as e:
        print(f"Error extracting text from PDF: {str(e)}")
        return None

def extract_text_from_docx(docx_path: str) -> Optional[str]:
    try:
        doc = docx.Document(docx_path)
        full_text = []
        for para in doc.paragraphs:
            if para.text.strip():
                full_text.append(para.text)
        return clean_text("\n".join(full_text))
    except Exception as e:
        print(f"Error extracting text from DOCX: {str(e)}")
        return None

def extract_text_from_txt(txt_path: str) -> Optional[str]:
    try:
        with open(txt_path, 'r', encoding='utf-8') as file:
            return clean_text(file.read())
    except Exception as e:
        print(f"Error extracting text from TXT: {str(e)}")
        return None

def extract_text_from_pptx(pptx_path: str) -> Optional[str]:
    try:
        prs = Presentation(pptx_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return clean_text("\n".join(text))
    except Exception as e:
        print(f"Error extracting text from PPTX: {str(e)}")
        return None

def extract_text_from_xlsx(xlsx_path: str) -> Optional[str]:
    try:
        xl = pd.ExcelFile(xlsx_path)
        text = []
        for sheet_name in xl.sheet_names:
            df = xl.parse(sheet_name)
            text.append(f"Sheet: {sheet_name}\n{df.to_string()}")
        return clean_text("\n\n".join(text))
    except Exception as e:
        print(f"Error extracting text from XLSX: {str(e)}")
        return None

def extract_text_from_csv(csv_path: str) -> Optional[str]:
    try:
        df = pd.read_csv(csv_path)
        return clean_text(df.to_string())
    except Exception as e:
        print(f"Error extracting text from CSV: {str(e)}")
        return None

def clean_text(text: str) -> str:
    # Remove excessive whitespace and line breaks
    text = re.sub(r'\s+', ' ', text).strip()
    # Remove non-printable characters
    text = ''.join(char for char in text if char.isprintable() or char in '\n\r\t')
    return text

def extract_text_from_file(file_path: str) -> Optional[str]:
    file_ext = os.path.splitext(file_path)[1].lower()

    if file_ext == '.pdf':
        return extract_text_from_pdf(file_path)
    elif file_ext == '.docx':
        return extract_text_from_docx(file_path)
    elif file_ext == '.txt':
        return extract_text_from_txt(file_path)
    elif file_ext == '.pptx':
        return extract_text_from_pptx(file_path)
    elif file_ext == '.xlsx':
        return extract_text_from_xlsx(file_path)
    elif file_ext == '.csv':
        return extract_text_from_csv(file_path)
    else:
        try:
            # Fallback to textract for other formats
            text = textract.process(file_path).decode('utf-8')
            return clean_text(text)
        except Exception as e:
            print(f"Error extracting text from file: {str(e)}")
            return None

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python documentProcessor.py <file_path>")
        sys.exit(1)

    file_path = sys.argv[1]
    content = extract_text_from_file(file_path)

    if content:
        print(content)
    else:
        print("Failed to extract content from document")
        sys.exit(1)