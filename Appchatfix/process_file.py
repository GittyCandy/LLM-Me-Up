import sys
import os
import json
from PyPDF2 import PdfReader
from docx import Document
import pptx
from openpyxl import load_workbook


def extract_text_from_pdf(file_path):
    text = ""
    with open(file_path, 'rb') as file:
        reader = PdfReader(file)
        num_pages = len(reader.pages)
        for page in reader.pages:
            text += page.extract_text() + "\n\n"
    return {"text": text.strip(), "pages": num_pages}


def extract_text_from_docx(file_path):
    doc = Document(file_path)
    text = "\n".join([para.text for para in doc.paragraphs if para.text])
    return {"text": text, "pages": len(doc.paragraphs) // 50 + 1}


def extract_text_from_pptx(file_path):
    prs = pptx.Presentation(file_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return {"text": "\n".join(text), "pages": len(prs.slides)}


def extract_text_from_xlsx(file_path):
    wb = load_workbook(file_path)
    text = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        for row in ws.iter_rows(values_only=True):
            text.append("\t".join([str(cell) if cell is not None else "" for cell in row]))
    return {"text": "\n".join(text), "pages": len(wb.sheetnames)}


def extract_text_from_txt(file_path):
    with open(file_path, 'r', encoding='utf-8') as file:
        text = file.read()
    return {"text": text, "pages": len(text.split('\n')) // 50 + 1}


def process_file(file_path):
    _, ext = os.path.splitext(file_path)
    ext = ext.lower()

    try:
        if ext == '.pdf':
            return extract_text_from_pdf(file_path)
        elif ext == '.docx':
            return extract_text_from_docx(file_path)
        elif ext == '.pptx':
            return extract_text_from_pptx(file_path)
        elif ext == '.xlsx':
            return extract_text_from_xlsx(file_path)
        elif ext == '.txt':
            return extract_text_from_txt(file_path)
        else:
            return {"error": "Unsupported file format"}
    except Exception as e:
        return {"error": str(e)}


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print(json.dumps({"error": "File path not provided"}))
        sys.exit(1)

    file_path = sys.argv[1]
    result = process_file(file_path)
    print(json.dumps(result))